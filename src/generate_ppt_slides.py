import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    print("Generating PPT Presentation Deck...")
    BASE_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    OUTPUT_PPT = os.path.join(BASE_DIR, "outputs", "Presentation_Deck.pptx")
    SUBMISSION_PPT = os.path.join(BASE_DIR, "Ananya_Shahi_Submission", "PPT_Slides", "Presentation_Deck.pptx")
    
    prs = Presentation()
    
    # Theme colors: Deep Blue and Slate Gray
    c_blue = RGBColor(11, 29, 58)
    c_slate = RGBColor(90, 107, 124)
    
    slides_data = [
        {
            "title": "Mutual Fund Analytics Platform",
            "subtitle": "Capstone Project I -- Bluestock Fintech\nAuthor: Ananya Shahi\nRole: Data Analyst Intern\nDate: June 11, 2026",
            "is_title": True
        },
        {
            "title": "Project Objectives",
            "bullets": [
                "Build a Python-based end-to-end ETL pipeline.",
                "Construct a normalized SQL database holding scheme history, investor logs, and risk metrics.",
                "Implement Value at Risk (VaR), CVaR, and Sector Concentration (HHI Index) algorithms.",
                "Develop a composite scorecard system to identify recommended schemes.",
                "Deploy an interactive web-based dashboard (Streamlit) for data visualizations."
            ]
        },
        {
            "title": "System Architecture",
            "bullets": [
                "ETL Layer: Pandas for handling missing daily records, reindexing dates, and data validation.",
                "Database Engine: SQLite (bluestock_mf.db) running a 6-table Star Schema.",
                "Analytics Engine: Automated returns compounding (CAGR), Sharpe ratios, and sector concentration indices.",
                "UI Layer: Streamlit (Python) + Plotly for interactive web visualizations."
            ]
        },
        {
            "title": "Database Design (Star Schema)",
            "bullets": [
                "dim_fund: Master table containing fund names, expense ratios, plans, risk classifications.",
                "fact_nav: Net asset value history log (64,320 rows).",
                "fact_transactions: Investor details (~32,000 processed rows).",
                "fact_performance: Trailing return percentages and calculated ratios.",
                "fact_portfolio: Holding weights per stock."
            ]
        },
        {
            "title": "Data Ingestion (ETL)",
            "bullets": [
                "Standardized transaction inputs (validating transaction types).",
                "Implemented forward-fill (ffill()) logic for scheme NAVs on weekends and public market holidays.",
                "Removed duplicates and ensured data integrity across all 10 datasets."
            ]
        },
        {
            "title": "Performance Scorecard Methodology",
            "bullets": [
                "Score rankings are computed using a weighted composite formula:",
                "3-Year Trailing returns (CAGR): 30% Weight",
                "Sharpe Ratio (Risk-Adjusted Return): 25% Weight",
                "Alpha (Performance relative to benchmark): 20% Weight",
                "Expense Ratio (Lower cost is better): 15% Weight",
                "Max Drawdown (Downside risk): 10% Weight"
            ]
        },
        {
            "title": "Top 5 Recommended Funds",
            "bullets": [
                "Ranks calculated using the weighted composite score:",
                "1. Kotak Flexicap Fund - Regular - Growth (Score: 0.7175)",
                "2. SBI Small Cap Fund - Regular Plan - Growth (Score: 0.7063)",
                "3. ICICI Pru Liquid Fund - Regular - Growth (Score: 0.7050)",
                "4. HDFC Short Term Debt Fund - Regular - Growth (Score: 0.7025)",
                "5. Kotak Emerging Equity Fund - Regular - Growth (Score: 0.6825)"
            ]
        },
        {
            "title": "Trailing Returns Analysis",
            "bullets": [
                "Liquid Debt funds showed steady, low-volatility returns.",
                "Flexicap and Small-cap equity schemes dominated the 3-Year and 5-Year CAGR.",
                "Kotak Flexicap Fund achieved a composite rank 1 owing to high Sharpe and low expense profile."
            ]
        },
        {
            "title": "Downside Risk Assessment (Value at Risk)",
            "bullets": [
                "Calculated 95% Historical Daily Value at Risk (VaR) and CVaR for each scheme.",
                "High-beta equity funds showed 95% Daily VaR limits between -1.5% to -2.4%.",
                "Low-risk debt schemes maintained daily VaR limits below -0.4%, proving safe liquidity profiles."
            ]
        },
        {
            "title": "Portfolio Sector HHI Index",
            "bullets": [
                "HHI < 1500 (Highly Diversified): Most mutual funds evaluated scored around 1200-1400, indicating high sector diversification.",
                "HHI > 2000 (Moderately Concentrated): Specific mid-cap schemes registered higher concentrations inside Finance and IT sectors."
            ]
        },
        {
            "title": "Interactive Streamlit Web App",
            "bullets": [
                "Built a Python Streamlit app (src/app.py) for live visualizations.",
                "Features interactive KPI metrics cards, risk vs. return scatter plots, state-wise investor volume bars, and sector pie charts.",
                "Equipped with dynamic sidebar filters for Fund House, Category, and Plan types."
            ]
        },
        {
            "title": "Project Conclusions",
            "bullets": [
                "Developed a clean, automated database architecture.",
                "Integrated mathematical return-risk scorings to support logical fund recommendation decisions.",
                "Created visual, interactive dashboards for stakeholders.",
                "Codebase is fully executable with clean comments and version-controlled on GitHub."
            ]
        }
    ]
    
    for item in slides_data:
        if item.get("is_title"):
            slide_layout = prs.slide_layouts[0] # Title layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = item["title"]
            subtitle.text = item["subtitle"]
            
            # Formats
            title.text_frame.paragraphs[0].font.color.rgb = c_blue
            title.text_frame.paragraphs[0].font.bold = True
        else:
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = item["title"]
            title.text_frame.paragraphs[0].font.color.rgb = c_blue
            
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for bullet in item["bullets"]:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)
                p.level = 0
                
    # Save files
    os.makedirs(os.path.dirname(OUTPUT_PPT), exist_ok=True)
    os.makedirs(os.path.dirname(SUBMISSION_PPT), exist_ok=True)
    
    prs.save(OUTPUT_PPT)
    prs.save(SUBMISSION_PPT)
    print("PowerPoint presentation generated successfully!")

if __name__ == "__main__":
    create_presentation()
